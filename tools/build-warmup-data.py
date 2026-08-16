# -*- coding: utf-8 -*-
"""把「雄麒道館｜健身班熱身標準化教案」PPT 轉成 index.html 內嵌的熱身站表資料集。

來源：D:\\雄麒道館課程計畫\\熱身教案PPT\\健身班熱身教案.pptx（不在 repo 內，只有要重新產生時才需要）
用法：python tools/build-warmup-data.py [pptx路徑]
產出：改寫 index.html 中 WARMUP-DATA-BEGIN / WARMUP-DATA-END 之間的內容。

這份 PPT 本質上是一套決策規則（人數→站數→年齡程度→器材→教練人力→安全准入），
系統已經握有人數、名單、年齡、級別、可出席教練，因此把它轉成可執行的資料而不是靜態文件。
"""

import json
import re
import sys
from pathlib import Path

from pptx import Presentation

SRC = Path(sys.argv[1] if len(sys.argv) > 1
           else r"D:\雄麒道館課程計畫\熱身教案PPT\健身班熱身教案_10_20_30_40人標準化優化版.pptx")
OUT = Path(__file__).resolve().parent.parent / "index.html"


def slide_title(s):
    for sh in s.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t and "雄麒道館" not in t and not t.isdigit():
                return t
    return ""


def tables(s):
    return [sh.table for sh in s.shapes if sh.has_table]


def rows(tbl, skip_header=True):
    out = []
    for i, r in enumerate(tbl.rows):
        cells = [c.text.strip().replace("\n", "") for c in r.cells]
        if skip_header and i == 0:
            continue
        if any(cells):
            out.append(cells)
    return out


def texts(s):
    return [sh.text_frame.text.strip() for sh in s.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()]


def pick(ts, prefix):
    for t in ts:
        if t.startswith(prefix):
            return t[len(prefix):].strip()
    return ""


def after(ts, label):
    """投影片把「工作 / 10-20秒」拆成兩個文字框，取 label 的下一個。"""
    for i, t in enumerate(ts):
        if t == label and i + 1 < len(ts):
            return ts[i + 1]
    return ""


prs = Presentation(str(SRC))
data = {
    "ver": "雄麒健身班熱身標準化教案",
    "cards": [], "layout": [], "roles": [], "timeline": [], "whistle": [],
    "coachCfg": [], "purposes": [], "ageRules": {}, "levelDef": {},
    "venue": [], "mixed": [], "risk": [], "preflight": {"checks": [], "fallback": []},
}

HIGH_RISK = {"中欄架", "高欄架", "翻滾墊", "攀爬／懸吊欄杆", "競速踢擊", "反應燈球", "螢幕顏色反應", "跳繩機"}

for s in prs.slides:
    title = slide_title(s)
    ts = texts(s)
    tb = tables(s)

    # ---------- 20 張器材卡 ----------
    if title.startswith("器材訓練方式｜"):
        name = title.split("｜", 1)[1].strip()
        card = {
            "n": name,
            "goal": pick(ts, "目的："),
            "ages": pick(ts, "適用："),
            "work": after(ts, "工作"),
            "sets": after(ts, "組數"),
            "rest": after(ts, "休息"),
            "short": pick(ts, "器材不足："),
            "safety": ts[-1] if ts and ts[-2:-1] == ["安全"] else "",
            "risk": 1 if name in HIGH_RISK else 0,
        }
        for t in tb:
            for r in rows(t, skip_header=False):
                if len(r) < 2:
                    continue
                k, v = r[0], r[1]
                if k == "初階":
                    card["A"] = v
                elif k == "中階":
                    card["B"] = v
                elif k == "高階":
                    card["C"] = v
                elif k.startswith("降階"):
                    card["step"] = v
                elif k.startswith("跆拳道"):
                    card["tkd"] = v
                elif k in ("10", "20", "30", "40"):
                    card.setdefault("cfg", {})[k] = v
        data["cards"].append(card)

    # ---------- 人數 → 站數公式 ----------
    elif "共用公式" in title:
        for r in rows(tb[0]):
            if r and r[0].endswith("人"):
                data["layout"].append({
                    "n": int(re.sub(r"\D", "", r[0])),
                    "st": r[1], "per": r[2], "rec": r[3] if len(r) > 3 else "",
                })

    # ---------- 5 人角色循環 ----------
    elif "角色循環" in title:
        ROLE = ("挑戰者", "下一棒", "模仿者", "小教練", "任務員")
        job = {}
        for i, t in enumerate(ts):          # 站位圖是「角色 / 說明」成對文字框
            if t in ROLE and i + 1 < len(ts) and ts[i + 1] not in ROLE:
                job.setdefault(t, ts[i + 1])
        for t in tb:                        # 表格再補「每輪只做一件事」
            for r in rows(t):
                if len(r) >= 2 and r[0] in ROLE:
                    job[r[0]] = job.get(r[0], r[1])
                    data.setdefault("_detail", {})[r[0]] = r[1]
        for r in ROLE:
            if r in job:
                data["roles"].append({"r": r, "job": job[r],
                                      "do": data.get("_detail", {}).get(r, "")})
        data.pop("_detail", None)
        for t in ts:
            if t.startswith("換人："):
                data["rotate"] = t

    # ---------- 每站 3 分鐘時間軸 + 哨音 ----------
    elif "每站3分鐘" in title:
        for r in rows(tb[0]):
            if len(r) >= 3 and ":" in r[0]:
                data["timeline"].append({"t": r[0], "what": r[1], "who": r[2]})
        for t in ts:                        # 哨音規則是一個多行文字框
            if "哨" in t and "：" in t:
                for line in t.split("\n"):
                    if "：" in line:
                        k, v = line.split("：", 1)
                        data["whistle"].append({"s": k.strip(), "m": v.strip()})

    # ---------- 人數 × 教練配置 ----------
    elif "教練配置" in title:
        for r in rows(tb[0]):
            if r and r[0].endswith("人"):
                data["coachCfg"].append({
                    "n": int(re.sub(r"\D", "", r[0])),
                    "low": r[1], "high": r[2], "rec": r[3] if len(r) > 3 else "",
                })

    # ---------- 先選目的 ----------
    elif "先選目的" in title:
        for r in rows(tb[0]):
            if len(r) >= 3:
                data["purposes"].append({"p": r[0], "do": r[1], "tkd": r[2]})

    # ---------- 年齡原則 ----------
    elif title.endswith("熱身原則"):
        age = title.replace("熱身原則", "").strip()
        items = []
        for r in rows(tb[0]):
            if len(r) >= 3:
                items.append({"k": r[0], "how": r[1], "load": r[2]})
        safety = ts[-1] if ts and ts[-2:-1] == ["安全"] else ""
        data["ageRules"][age] = {"items": items, "safety": safety}

    # ---------- Age × Level 定義 ----------
    elif "Age × Level" in title or "年齡 × 程度" in title or "年齡不等於程度" in title:
        for t in ts:
            for line in t.split("\n"):
                m = re.match(r"^(初階|中階|高階)：(.+)$", line.strip())
                if m:
                    data["levelDef"][m.group(1)] = m.group(2)
        for t in tb:                        # Age × Level 矩陣
            hdr = [c.text.strip() for c in t.rows[0].cells]
            for r in list(t.rows)[1:]:
                cs = [c.text.strip().replace("\n", "／") for c in r.cells]
                if cs[0]:
                    data.setdefault("ageLevel", {})[cs[0]] = {
                        hdr[i]: cs[i] for i in range(1, len(cs)) if hdr[i]}

    # ---------- 器材附錄索引（館方自己的四類分法） ----------
    elif "器材附錄索引" in title:
        pairs = [t for t in ts if t not in ("器材附錄索引",) and "雄麒道館" not in t and not t.isdigit()]
        i = 0
        while i + 1 < len(pairs):
            cat, items = pairs[i], pairs[i + 1]
            if "、" in items:
                data.setdefault("index", []).append(
                    {"c": cat, "eq": [x.strip() for x in items.split("、") if x.strip()]})
                i += 2
            else:
                i += 1

    # ---------- 小場地 ----------
    elif "小場地版本" in title:
        for r in rows(tb[0]):
            if len(r) >= 4:
                data["venue"].append({"c": r[0], "j": r[1], "can": r[2], "no": r[3]})

    # ---------- 混齡混程度 ----------
    elif "混齡混程度" in title:
        for r in rows(tb[0]):
            if len(r) >= 4:
                data["mixed"].append({"eq": r[0], "A": r[1], "B": r[2], "C": r[3]})

    # ---------- 高風險站准入 ----------
    elif "高風險站准入" in title:
        for r in rows(tb[0]):
            if len(r) >= 4:
                data["risk"].append({"eq": r[0], "s": r[1], "l": r[2], "no": r[3]})

    # ---------- 電子設備 PRE-FLIGHT ----------
    elif "PRE-FLIGHT" in title or "課前檢查" in title:
        for t in tb:
            for r in rows(t, skip_header=False):
                if len(r) < 2:
                    continue
                if r[0] in ("設備", "確認"):
                    continue
                if r[1] in ("足夠", "已連線", "已登入", "已測試", "已設定"):
                    data["preflight"]["checks"].append({"k": r[0], "v": r[1]})
                else:
                    data["preflight"]["fallback"].append({"eq": r[0], "alt": r[1]})
        for t in ts:
            if "30秒" in t:
                data["preflight"]["rule"] = t

payload = ("/* 雄麒健身班熱身標準化教案｜由 tools/build-warmup-data.py 自動產生，請勿手改 */\n"
           "const WSOP=" + json.dumps(data, ensure_ascii=False, separators=(",", ":")) + ";")

html = OUT.read_text(encoding="utf-8")
B, E = "/* WARMUP-DATA-BEGIN */", "/* WARMUP-DATA-END */"
if B not in html:
    print("index.html 找不到 WARMUP-DATA-BEGIN 標記")
    sys.exit(1)
html = html[:html.index(B) + len(B)] + "\n" + payload + "\n" + html[html.index(E):]
OUT.write_text(html, encoding="utf-8")

print(f"器材卡 {len(data['cards'])}｜人數公式 {len(data['layout'])}｜角色 {len(data['roles'])}｜"
      f"時間軸 {len(data['timeline'])}｜哨音 {len(data['whistle'])}｜教練配置 {len(data['coachCfg'])}")
print(f"目的 {len(data['purposes'])}｜年齡原則 {len(data['ageRules'])}｜程度定義 {len(data['levelDef'])}｜"
      f"場地 {len(data['venue'])}｜混齡 {len(data['mixed'])}｜高風險 {len(data['risk'])}｜"
      f"PRE-FLIGHT {len(data['preflight']['checks'])}+{len(data['preflight']['fallback'])}")
print(f"資料大小 {len(payload)/1024:.1f} KB")
